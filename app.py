import streamlit as st
from pdf2docx import Converter
import pytesseract
from pdf2image import convert_from_bytes, convert_from_path
from docx import Document
from pptx import Presentation
from PIL import Image
import os
import subprocess
import zipfile
import io
import platform # Ä°ÅŸletim sistemini anlamak iÃ§in
import shutil   # Linux'ta komut kontrolÃ¼ iÃ§in
from PyPDF2 import PdfMerger, PdfReader, PdfWriter

# NOT: pythoncom ve docx2pdf kÃ¼tÃ¼phaneleri Linux'ta Ã§alÄ±ÅŸmadÄ±ÄŸÄ± iÃ§in kaldÄ±rÄ±ldÄ±.
# ArtÄ±k LibreOffice kullanÄ±yoruz, onlara ihtiyacÄ±mÄ±z yok.

# ========================================================
#                  AYARLAR VE SABÄ°TLER (CROSS-PLATFORM)
# ========================================================

# Ä°ÅŸletim Sistemi KontrolÃ¼
if platform.system() == "Windows":
    # SENÄ°N BÄ°LGÄ°SAYARINDAKÄ° YOLLAR (Windows)
    TESSERACT_PATH = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    POPPLER_PATH = r'C:\Program Files\poppler-24.02.0\Library\bin'
    LIBREOFFICE_PATH = r'C:\Program Files\LibreOffice\program\soffice.exe'
else:
    # STREAMLIT CLOUD (LINUX) YOLLARI
    TESSERACT_PATH = "tesseract"
    POPPLER_PATH = None 
    LIBREOFFICE_PATH = "soffice"

# Tesseract yolunu ata
pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH

st.set_page_config(page_title="Type-Pdf", layout="centered",)
st.title("Type-Pdf")

# ========================================================
#                  MENÃœ SÄ°STEMÄ°
# ========================================================

st.sidebar.header(" MenÃ¼ YÃ¶netimi")

kategori = st.sidebar.selectbox(
    "Kategori SeÃ§iniz:",
    ["DÃ¶nÃ¼ÅŸtÃ¼rme Ä°ÅŸlemleri (Converter)", 
     "PDF AraÃ§larÄ± (Tools)",
     "Sistem ve YardÄ±m (System Info)"]
)

secim = ""

if kategori == "DÃ¶nÃ¼ÅŸtÃ¼rme Ä°ÅŸlemleri (Converter)":
    st.sidebar.subheader(" DÃ¶nÃ¼ÅŸtÃ¼rme ModlarÄ±")
    secim = st.sidebar.radio("Ä°ÅŸlem SeÃ§iniz:", 
        ["PDF -> Word (Metin)", 
         "Word -> PDF (LibreOffice)",
         "Word -> JPG (LibreOffice)",
         "PDF -> RTF (Zengin Metin)",
         "RTF -> PDF",
         "JPG -> PDF (Resimden PDF)",
         "JPG -> Word (OCR)",
         "PowerPoint -> PDF (LibreOffice)",
         "PDF -> PowerPoint (Sunum)",
         "OCR: TaranmÄ±ÅŸ PDF -> Word"])

elif kategori == "PDF AraÃ§larÄ± (Tools)":
    st.sidebar.subheader(" PDF AraÃ§larÄ±")
    secim = st.sidebar.radio("AraÃ§ SeÃ§iniz:", 
        ["PDF BirleÅŸtir (Merge)", 
         "PDF BÃ¶l (Split)", 
         "PDF SÄ±kÄ±ÅŸtÄ±rma (Optimizer)", 
         "PDF Åžifrele (Lock)",
         "PDF Metadata Temizle (Privacy)"])

elif kategori == "Sistem ve YardÄ±m (System Info)":
    secim = "Sistem Durumu"

# ========================================================
#                  MODÃœL 0: SÄ°STEM SAÄžLIK KONTROLÃœ
# ========================================================

if secim == "Sistem Durumu":
    st.header(" Sistem SaÄŸlÄ±k ve BaÄŸÄ±mlÄ±lÄ±k KontrolÃ¼")
    st.info("Bu panel, uygulamanÄ±n Ã§alÄ±ÅŸmasÄ± iÃ§in gerekli harici araÃ§larÄ±n durumunu gÃ¶sterir.")
    
    col1, col2, col3 = st.columns(3)
    
    # Kontrol Fonksiyonu (Windows ve Linux uyumlu)
    def check_tool(path, name_linux):
        if platform.system() == "Windows":
            return os.path.exists(path)
        else:
            return shutil.which(name_linux) is not None

    # 1. Tesseract KontrolÃ¼
    with col1:
        if check_tool(TESSERACT_PATH, "tesseract"):
            st.success("Tesseract OCR")
            st.caption(f"Durum: **Aktif**")
        else:
            st.error("Tesseract OCR")
            st.caption("Durum: **BulunamadÄ±!**")

    # 2. Poppler KontrolÃ¼
    with col2:
        if platform.system() == "Windows":
            status = os.path.exists(POPPLER_PATH)
        else:
            status = shutil.which("pdftoppm") is not None
            
        if status:
            st.success("Poppler Utils")
            st.caption(f"Durum: **Aktif**")
        else:
            st.error("Poppler Utils")
            st.caption("Durum: **BulunamadÄ±!**")

    # 3. LibreOffice KontrolÃ¼
    with col3:
        if check_tool(LIBREOFFICE_PATH, "soffice"):
            st.success("LibreOffice")
            st.caption(f"Durum: **Aktif**")
        else:
            st.error("LibreOffice")
            st.caption("Durum: **BulunamadÄ±!**")

    st.divider()
    st.subheader("ðŸ‘¨â€ðŸ’» Proje HakkÄ±nda")
    st.write(f"""
    **Ã‡alÄ±ÅŸma OrtamÄ±:** {platform.system()}
    
    Bu proje, **Bilgisayar MÃ¼hendisliÄŸi UygulamalarÄ±** dersi kapsamÄ±nda geliÅŸtirilmiÅŸtir.
    """)

# ========================================================
#                  MODÃœL 1: PDF ARAÃ‡LARI
# ========================================================

elif secim == "PDF BirleÅŸtir (Merge)":
    st.header(" PDF DosyalarÄ±nÄ± BirleÅŸtir")
    uploaded_pdfs = st.file_uploader("PDF'leri SeÃ§in", type="pdf", accept_multiple_files=True, key="merge")
    if uploaded_pdfs and st.button("BirleÅŸtir"):
        with st.spinner('BirleÅŸtiriliyor...'):
            try:
                merger = PdfMerger()
                for pdf in uploaded_pdfs: merger.append(pdf)
                buf = io.BytesIO()
                merger.write(buf)
                merger.close()
                st.download_button("Ä°ndir", buf.getvalue(), "birlestirilmis.pdf", "application/pdf")
                st.success("Bitti!")
            except Exception as e: st.error(f"Hata: {e}")

elif secim == "PDF BÃ¶l (Split)":
    st.header(" PDF DosyasÄ±nÄ± BÃ¶l")
    up_split = st.file_uploader("PDF YÃ¼kle", type="pdf", key="split")
    if up_split:
        reader = PdfReader(up_split)
        total = len(reader.pages)
        st.write(f"Toplam Sayfa: {total}")
        c1, c2 = st.columns(2)
        start = c1.number_input("BaÅŸlangÄ±Ã§", 1, total, 1)
        end = c2.number_input("BitiÅŸ", 1, total, total)
        if st.button("BÃ¶l ve Ä°ndir"):
            if start > end: st.error("HatalÄ± aralÄ±k.")
            else:
                writer = PdfWriter()
                for i in range(start-1, end): writer.add_page(reader.pages[i])
                buf = io.BytesIO()
                writer.write(buf)
                st.download_button("Ä°ndir", buf.getvalue(), f"bolunmus_{start}-{end}.pdf", "application/pdf")
                st.success("Bitti!")

elif secim == "PDF SÄ±kÄ±ÅŸtÄ±rma (Optimizer)":
    st.header(" PDF Boyut KÃ¼Ã§Ã¼ltme")
    st.info("PDF iÃ§indeki gereksiz boÅŸluklarÄ± ve akÄ±ÅŸlarÄ± temizler.")
    up_opt = st.file_uploader("PDF YÃ¼kle", type="pdf", key="compress")
    if up_opt:
        original_size = up_opt.size / 1024
        st.write(f" **Orijinal Boyut:** {original_size:.2f} KB")
        if st.button("SÄ±kÄ±ÅŸtÄ±r"):
            with st.spinner("Optimize ediliyor..."):
                try:
                    reader = PdfReader(up_opt)
                    writer = PdfWriter()
                    for page in reader.pages:
                        page.compress_content_streams() 
                        writer.add_page(page)
                    
                    buf = io.BytesIO()
                    writer.write(buf)
                    new_size = buf.getbuffer().nbytes / 1024
                    ratio = ((original_size - new_size) / original_size) * 100
                    st.write(f"ðŸ“¦ **Yeni Boyut:** {new_size:.2f} KB")
                    if new_size < original_size:
                        st.success(f"BaÅŸarÄ±lÄ±! %{ratio:.1f} oranÄ±nda sÄ±kÄ±ÅŸtÄ±.")
                    else: st.info("Dosya zaten optimize edilmiÅŸ.")
                    st.download_button("Ä°ndir", buf.getvalue(), "optimize.pdf", "application/pdf")
                except Exception as e: st.error(f"Hata: {e}")

elif secim == "PDF Åžifrele (Lock)":
    st.header(" PDF Åžifreleme")
    up_lock = st.file_uploader("PDF YÃ¼kle", type="pdf", key="lock")
    if up_lock:
        pwd = st.text_input("Åžifre", type="password")
        if st.button("Kilitle"):
            if pwd:
                reader = PdfReader(up_lock)
                writer = PdfWriter()
                for p in reader.pages: writer.add_page(p)
                writer.encrypt(pwd)
                buf = io.BytesIO()
                writer.write(buf)
                st.download_button("Ä°ndir", buf.getvalue(), "sifreli.pdf", "application/pdf")
                st.success("Kilitlendi.")
            else: st.warning("Åžifre giriniz.")

elif secim == "PDF Metadata Temizle (Privacy)":
    st.header(" PDF Metadata Temizleme")
    up_meta = st.file_uploader("PDF YÃ¼kle", type="pdf", key="meta")
    if up_meta:
        reader = PdfReader(up_meta)
        st.json(reader.metadata)
        if st.button("Temizle"):
            writer = PdfWriter()
            for p in reader.pages: writer.add_page(p)
            buf = io.BytesIO()
            writer.write(buf)
            st.download_button("Ä°ndir", buf.getvalue(), "temiz.pdf", "application/pdf")
            st.success("Metadata silindi.")

# ========================================================
#                  MODÃœL 2: DÃ–NÃœÅžTÃœRME
# ========================================================

elif secim == "PDF -> Word (Metin)":
    st.header(" PDF -> Word")
    up = st.file_uploader("PDF", type="pdf", key="p2w")
    if up and st.button("Ã‡evir"):
        with st.spinner('...'):
            try:
                with open("t.pdf", "wb") as f: f.write(up.getbuffer())
                cv = Converter("t.pdf")
                cv.convert("o.docx")
                cv.close()
                with open("o.docx", "rb") as f: st.download_button("Ä°ndir", f, "d.docx")
                st.success("Ok")
            except Exception as e: st.error(e)

elif secim == "Word -> PDF (LibreOffice)":
    st.header(" Word -> PDF")
    up = st.file_uploader("Word", type="docx", key="w2p")
    if up and st.button("Ã‡evir"):
        with st.spinner('LibreOffice...'):
            try:
                base = os.getcwd()
                inp = os.path.join(base, "t.docx")
                with open(inp, "wb") as f: f.write(up.getbuffer())
                cmd = [LIBREOFFICE_PATH, '--headless', '--convert-to', 'pdf', inp, '--outdir', base]
                subprocess.run(cmd, capture_output=True)
                if os.path.exists("t.pdf"):
                    with open("t.pdf", "rb") as f: st.download_button("Ä°ndir", f, "d.pdf")
                    st.success("Ok")
                    os.remove(inp); os.remove("t.pdf")
                else: st.error("Hata: LibreOffice yolu doÄŸru mu?")
            except Exception as e: st.error(e)

elif secim == "Word -> JPG (LibreOffice)":
    st.header(" Word -> JPG")
    up = st.file_uploader("Word", type="docx", key="w2j")
    if up and st.button("Ã‡evir"):
        with st.spinner('...'):
            try:
                base = os.getcwd()
                inp = os.path.join(base, "ti.docx")
                with open(inp, "wb") as f: f.write(up.getbuffer())
                cmd = [LIBREOFFICE_PATH, '--headless', '--convert-to', 'pdf', inp, '--outdir', base]
                subprocess.run(cmd, capture_output=True)
                tpdf = os.path.join(base, "ti.pdf")
                if os.path.exists(tpdf):
                    imgs = convert_from_path(tpdf, poppler_path=POPPLER_PATH)
                    buf = io.BytesIO()
                    with zipfile.ZipFile(buf, "w") as z:
                        for i, im in enumerate(imgs):
                            b = io.BytesIO()
                            im.save(b, 'JPEG')
                            z.writestr(f"p_{i+1}.jpg", b.getvalue())
                    st.download_button("ZIP Ä°ndir", buf.getvalue(), "img.zip", "application/zip")
                    st.success("Ok")
                    os.remove(inp); os.remove(tpdf)
                else: st.error("Hata")
            except Exception as e: st.error(e)

elif secim == "PDF -> RTF (Zengin Metin)":
    st.header(" PDF -> RTF")
    up = st.file_uploader("PDF", type="pdf", key="p2rtf")
    if up and st.button("Ã‡evir"):
        with st.spinner('...'):
            try:
                base = os.getcwd()
                with open("tc.pdf", "wb") as f: f.write(up.getbuffer())
                cv = Converter("tc.pdf")
                cv.convert("tc.docx")
                cv.close()
                cmd = [LIBREOFFICE_PATH, '--headless', '--convert-to', 'rtf', "tc.docx", '--outdir', base]
                subprocess.run(cmd, capture_output=True)
                if os.path.exists("tc.rtf"):
                    with open("tc.rtf", "rb") as f: st.download_button("Ä°ndir", f, "b.rtf")
                    st.success("Ok")
                else: st.error("Hata")
            except Exception as e: st.error(e)

elif secim == "RTF -> PDF":
    st.header(" RTF -> PDF")
    up = st.file_uploader("RTF", type="rtf", key="r2p")
    if up and st.button("Ã‡evir"):
        with st.spinner('...'):
            try:
                base = os.getcwd()
                inp = "temp.rtf"
                with open(inp, "wb") as f: f.write(up.getbuffer())
                cmd = [LIBREOFFICE_PATH, '--headless', '--convert-to', 'pdf', inp, '--outdir', base]
                subprocess.run(cmd, capture_output=True)
                if os.path.exists("temp.pdf"):
                    with open("temp.pdf", "rb") as f: st.download_button("Ä°ndir", f, "b.pdf")
                    st.success("Ok")
                else: st.error("Hata")
            except Exception as e: st.error(e)

elif secim == "JPG -> PDF (Resimden PDF)":
    st.header(" JPG -> PDF")
    ups = st.file_uploader("JPG", type=["jpg","png"], accept_multiple_files=True, key="j2p")
    if ups and st.button("Ã‡evir"):
        try:
            lst = []
            for u in ups:
                i = Image.open(u).convert('RGB')
                lst.append(i)
            if lst:
                lst[0].save("m.pdf", save_all=True, append_images=lst[1:])
                with open("m.pdf", "rb") as f: st.download_button("Ä°ndir", f, "m.pdf")
                st.success("Ok")
        except Exception as e: st.error(e)

elif secim == "JPG -> Word (OCR)":
    st.header(" JPG -> Word")
    up = st.file_uploader("Resim", type=["jpg","png"], key="j2w")
    if up and st.button("Ã‡evir"):
        try:
            txt = pytesseract.image_to_string(Image.open(up), lang='tur')
            d = Document()
            d.add_paragraph(txt)
            d.save("o.docx")
            with open("o.docx", "rb") as f: st.download_button("Ä°ndir", f, "o.docx")
            st.success("Ok")
        except Exception as e: st.error(e)

elif secim == "PowerPoint -> PDF (LibreOffice)":
    st.header(" PPT -> PDF")
    up = st.file_uploader("PPT", type="ppt", key="pp2p")
    if up and st.button("Ã‡evir"):
        try:
            base = os.getcwd()
            inp = os.path.join(base, "t.pptx")
            with open(inp, "wb") as f: f.write(up.getbuffer())
            subprocess.run([LIBREOFFICE_PATH, '--headless', '--convert-to', 'pdf', inp, '--outdir', base])
            if os.path.exists("t.pdf"):
                with open("t.pdf", "rb") as f: st.download_button("Ä°ndir", f, "s.pdf")
                st.success("Ok")
        except Exception as e: st.error(e)

elif secim == "PDF -> PowerPoint (Sunum)":
    st.header(" PDF -> PPTX")
    up = st.file_uploader("PDF", type="pdf", key="p2pp")
    if up and st.button("Ã‡evir"):
        try:
            imgs = convert_from_bytes(up.read(), poppler_path=POPPLER_PATH)
            prs = Presentation()
            for i, im in enumerate(imgs):
                im.save(f"t_{i}.jpg")
                s = prs.slides.add_slide(prs.slide_layouts[6])
                s.shapes.add_picture(f"t_{i}.jpg", 0, 0, height=prs.slide_height)
                os.remove(f"t_{i}.jpg")
            prs.save("c.pptx")
            with open("c.pptx", "rb") as f: st.download_button("Ä°ndir", f, "s.pptx")
            st.success("Ok")
        except Exception as e: st.error(e)

elif secim == "OCR: TaranmÄ±ÅŸ PDF -> Word":
    st.header(" OCR PDF")
    up = st.file_uploader("PDF", type="pdf", key="ocr")
    if up and st.button("Ã‡evir"):
        try:
            imgs = convert_from_bytes(up.read(), poppler_path=POPPLER_PATH)
            d = Document()
            for im in imgs:
                t = pytesseract.image_to_string(im, lang='tur')
                d.add_paragraph(t)
                d.add_page_break()
            d.save("ocr.docx")
            with open("ocr.docx", "rb") as f: st.download_button("Ä°ndir", f, "ocr.docx")
            st.success("Ok")
        except Exception as e: st.error(e)